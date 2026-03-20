#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT 美化脚本 - 修复排版问题，统一字体和样式
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import sys

# 设计规范
DESIGN_SPECS = {
    'background_color': RGBColor(245, 245, 245),  # #F5F5F5
    'primary_color': RGBColor(30, 58, 138),        # #1E3A8A
    'text_color': RGBColor(31, 41, 55),            # #1F2937
    'title_font': 'Microsoft YaHei',
    'title_size': Pt(24),
    'body_font': 'Microsoft YaHei',
    'body_size': Pt(14),
}

def set_font(shape, font_name='Microsoft YaHei', font_size=Pt(14), bold=False, color=None):
    """设置形状的字体"""
    if hasattr(shape, 'text_frame') and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = font_size
                run.font.bold = bold
                if color:
                    run.font.color.rgb = color

def fix_text_overlap(shape, slide_width, slide_height):
    """修复文字重叠问题"""
    if hasattr(shape, 'text_frame') and shape.text_frame:
        # 调整文本框内边距
        tf = shape.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        # 检查是否超出边界 - 添加安全检查
        min_size = Inches(0.5)
        max_width = slide_width - shape.left - Inches(0.2)
        max_height = slide_height - shape.top - Inches(0.2)
        
        if max_width > min_size and shape.width > max_width:
            shape.width = max_width
        if max_height > min_size and shape.height > max_height:
            shape.height = max_height

def fix_table(table):
    """修复表格错位问题"""
    if table:
        # 统一单元格内边距
        for row in table.rows:
            for cell in row.cells:
                tf = cell.text_frame
                tf.margin_left = Inches(0.05)
                tf.margin_right = Inches(0.05)
                tf.margin_top = Inches(0.05)
                tf.margin_bottom = Inches(0.05)
                tf.word_wrap = True
                
                # 设置单元格字体
                for paragraph in tf.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.name = 'Microsoft YaHei'
                        run.font.size = Pt(12)

def beautify_slide(slide, slide_width, slide_height):
    """美化单页幻灯片"""
    # 设置背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DESIGN_SPECS['background_color']
    
    # 遍历所有形状
    for shape in slide.shapes:
        # 跳过某些特殊形状
        if hasattr(shape, 'shape_type') and shape.shape_type in [MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM]:
            continue
            
        # 处理文本框和形状
        if hasattr(shape, 'text_frame'):
            # 修复文字重叠
            try:
                fix_text_overlap(shape, slide_width, slide_height)
            except Exception as e:
                pass  # 忽略调整失败
            
            # 判断是标题还是正文
            is_title = False
            if hasattr(shape, 'name') and ('标题' in str(shape.name).lower() or 'title' in str(shape.name).lower()):
                is_title = True
            elif shape.text_frame.paragraphs and len(shape.text_frame.paragraphs) > 0:
                first_run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                if first_run and first_run.font.size and first_run.font.size > Pt(18):
                    is_title = True
            
            # 设置字体
            try:
                if is_title:
                    set_font(shape, DESIGN_SPECS['title_font'], DESIGN_SPECS['title_size'], True, DESIGN_SPECS['text_color'])
                else:
                    set_font(shape, DESIGN_SPECS['body_font'], DESIGN_SPECS['body_size'], False, DESIGN_SPECS['text_color'])
            except Exception as e:
                pass  # 忽略字体设置失败
        
        # 处理表格
        if hasattr(shape, 'table') and shape.table:
            try:
                fix_table(shape.table)
            except Exception as e:
                pass  # 忽略表格修复失败

def beautify_ppt(input_path, output_path=None):
    """美化 PPT 文件"""
    if output_path is None:
        base, ext = os.path.splitext(input_path)
        output_path = f"{base}-fixed{ext}"
    
    print(f"Processing: {input_path}")
    
    try:
        prs = Presentation(input_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 美化每一页
        for i, slide in enumerate(prs.slides, 1):
            print(f"  - Slide {i}/{len(prs.slides)}")
            try:
                beautify_slide(slide, slide_width, slide_height)
            except Exception as e:
                print(f"    Warning: Slide {i} error: {str(e)}")
        
        # 保存文件
        prs.save(output_path)
        print(f"  [OK] Saved to: {output_path}")
        return True
        
    except Exception as e:
        print(f"  [ERROR] Failed: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

def find_ppt_files():
    """查找所有 PPT 文件"""
    workspace_alphasage = r"C:\Users\Administrator\.openclaw\workspace-alphasage"
    workspace = r"C:\Users\Administrator\.openclaw\workspace"
    
    ppt_files = []
    
    # 查找 workspace 中的文件
    if os.path.exists(workspace):
        for f in os.listdir(workspace):
            if f.endswith('.pptx') and '端云协同' in f:
                ppt_files.append(os.path.join(workspace, f))
    
    # 查找 workspace-alphasage 中的文件
    if os.path.exists(workspace_alphasage):
        for f in os.listdir(workspace_alphasage):
            if f.endswith('.pptx') and 'AI-Agents' in f and 'fixed' not in f:
                ppt_files.append(os.path.join(workspace_alphasage, f))
    
    return ppt_files

def main():
    ppt_files = find_ppt_files()
    
    print("=" * 60)
    print("PPT Beautification Script - Starting")
    print("=" * 60)
    print(f"Found {len(ppt_files)} files:")
    for f in ppt_files:
        print(f"  - {f}")
    print(f"\nDesign Specs:")
    print(f"  - Title: {DESIGN_SPECS['title_font']} Bold, {DESIGN_SPECS['title_size']}")
    print(f"  - Body: {DESIGN_SPECS['body_font']} Regular, {DESIGN_SPECS['body_size']}")
    print(f"  - Background: #F5F5F5, Primary: #1E3A8A, Text: #1F2937")
    print("=" * 60)
    
    success_count = 0
    for ppt_file in ppt_files:
        if beautify_ppt(ppt_file):
            success_count += 1
    
    print("=" * 60)
    print(f"Completed: {success_count}/{len(ppt_files)} files successful")
    print("=" * 60)

if __name__ == "__main__":
    main()
