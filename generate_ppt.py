#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
端云协同 AI Agent 在水利设计院的创新应用
5 页 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 颜色定义
BG_COLOR = RGBColor(245, 245, 245)  # 浅灰白背景
PRIMARY_COLOR = RGBColor(30, 58, 138)  # 深蓝主色
ACCENT_COLOR = RGBColor(59, 130, 246)  # 亮蓝强调色
TEXT_COLOR = RGBColor(31, 41, 55)  # 深灰文字
WHITE = RGBColor(255, 255, 255)

# 辅助颜色
LIGHT_BLUE = RGBColor(227, 242, 253)
LIGHT_GREEN = RGBColor(232, 245, 233)
LIGHT_ORANGE = RGBColor(255, 243, 224)
LIGHT_RED = RGBColor(255, 235, 238)
LIGHT_PURPLE = RGBColor(232, 234, 246)

def set_slide_background(slide, prs):
    """设置幻灯片背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title_slide(slide, title_text):
    """添加标题"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.name = '微软雅黑'
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    return title_box

def add_bullet_points(slide, left, top, width, height, points, font_size=16):
    """添加项目符号列表"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            para = frame.paragraphs[0]
        else:
            para = frame.add_paragraph()
        para.text = point
        para.font.name = '微软雅黑'
        para.font.size = Pt(font_size)
        para.font.color.rgb = TEXT_COLOR
        para.level = 0
    return textbox

def add_table(slide, left, top, width, height, rows, cols, data, header_color=PRIMARY_COLOR):
    """添加表格"""
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 设置列宽
    col_width = int(width / cols)
    for col in range(cols):
        table.columns[col].width = col_width
    
    # 填充数据
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            
            # 设置单元格文本格式
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = '微软雅黑'
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = TEXT_COLOR
                
                if row_idx == 0:  # 表头
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_color
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
    
    return table

def add_architecture_diagram(slide):
    """添加三层架构图"""
    # 云端
    cloud = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.2), Inches(8), Inches(1.5))
    cloud.fill.solid()
    cloud.fill.fore_color.rgb = LIGHT_BLUE
    cloud.line.color.rgb = PRIMARY_COLOR
    cloud.line.width = Pt(2)
    
    cloud_text = "云端（省院信息中心）\n"
    cloud_text += "• 国产算力集群（华为昇腾/寒武纪）\n"
    cloud_text += "• 私有化大模型部署（Qwen/GLM）\n"
    cloud_text += "• 数据不出域，统一训练与推理"
    
    cloud_frame = cloud.text_frame
    cloud_frame.clear()
    para = cloud_frame.paragraphs[0]
    para.text = cloud_text
    para.font.name = '微软雅黑'
    para.font.size = Pt(14)
    para.font.color.rgb = TEXT_COLOR
    para.alignment = PP_ALIGN.LEFT
    
    # 箭头 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.75), Inches(2.8), Inches(0.5), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = ACCENT_COLOR
    arrow1.line.color.rgb = ACCENT_COLOR
    
    # 边侧
    edge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3.3), Inches(8), Inches(1.5))
    edge.fill.solid()
    edge.fill.fore_color.rgb = LIGHT_GREEN
    edge.line.color.rgb = PRIMARY_COLOR
    edge.line.width = Pt(2)
    
    edge_text = "边侧（客户数据中心/水库管理处）\n"
    edge_text += "• 边缘计算节点\n"
    edge_text += "• 区域数据汇聚与预处理\n"
    edge_text += "• 轻量级 Agent 协同调度"
    
    edge_frame = edge.text_frame
    edge_frame.clear()
    para = edge_frame.paragraphs[0]
    para.text = edge_text
    para.font.name = '微软雅黑'
    para.font.size = Pt(14)
    para.font.color.rgb = TEXT_COLOR
    para.alignment = PP_ALIGN.LEFT
    
    # 箭头 2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.75), Inches(4.9), Inches(0.5), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = ACCENT_COLOR
    arrow2.line.color.rgb = ACCENT_COLOR
    
    # 端侧
    device = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.4), Inches(8), Inches(1.5))
    device.fill.solid()
    device.fill.fore_color.rgb = LIGHT_ORANGE
    device.line.color.rgb = PRIMARY_COLOR
    device.line.width = Pt(2)
    
    device_text = "端侧（闸门/监测设备/传感器）\n"
    device_text += "• 端侧 Agent 嵌入闸控系统\n"
    device_text += "• 实时决策与应急响应\n"
    device_text += "• 离线可用，断网仍可运行核心功能"
    
    device_frame = device.text_frame
    device_frame.clear()
    para = device_frame.paragraphs[0]
    para.text = device_text
    para.font.name = '微软雅黑'
    para.font.size = Pt(14)
    para.font.color.rgb = TEXT_COLOR
    para.alignment = PP_ALIGN.LEFT

def add_flow_diagram(slide):
    """添加数据通路流程图"""
    positions = [
        (Inches(0.8), Inches(1.3), "传感器层（水位、流量、闸门开度）"),
        (Inches(0.8), Inches(2.0), "端侧 Agent（嵌入式设备，如 RK3588/Jetson）"),
        (Inches(0.8), Inches(2.7), "智能分析 + 建议生成（推荐开度、时机、流量）"),
        (Inches(0.8), Inches(3.4), "⚠️ 人工确认（值班人员审核批准）"),
        (Inches(0.8), Inches(4.1), "执行指令（闸控系统动作）"),
        (Inches(0.8), Inches(4.8), "边侧同步（水库数据中心）"),
        (Inches(0.8), Inches(5.5), "云端监控（省院信息中心）"),
    ]
    
    for i, (left, top, text) in enumerate(positions):
        # 添加矩形框
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(8.4), Inches(0.5))
        box.fill.solid()
        if "⚠️" in text:
            box.fill.fore_color.rgb = LIGHT_RED
            box.line.color.rgb = RGBColor(211, 47, 47)
        else:
            box.fill.fore_color.rgb = LIGHT_BLUE
            box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(1.5)
        
        # 添加文本
        box_frame = box.text_frame
        box_frame.clear()
        para = box_frame.paragraphs[0]
        para.text = text
        para.font.name = '微软雅黑'
        para.font.size = Pt(13)
        para.font.color.rgb = TEXT_COLOR
        para.alignment = PP_ALIGN.CENTER
        
        # 添加箭头（除了最后一个）
        if i < len(positions) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.5), top + Inches(0.55), Inches(0.3), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_COLOR
            arrow.line.color.rgb = ACCENT_COLOR

def create_ppt():
    """创建 PPT"""
    prs = Presentation()
    
    # 设置幻灯片尺寸为 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========== 第 1 页：核心框架设计 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_background(slide1, prs)
    
    add_title_slide(slide1, "核心框架设计 —— 安全优先的端云协同架构")
    add_architecture_diagram(slide1)
    
    # 安全设计要点
    safety_points = [
        "✅ 数据不出域：敏感数据（水文、地理信息）留在本地",
        "✅ 国产算力：适配华为昇腾、寒武纪等国产芯片",
        "✅ 分层授权：云端训练、边侧推理、端侧执行",
        "✅ 审计追溯：全链路操作日志，符合等保 2.0 要求"
    ]
    add_bullet_points(slide1, Inches(0.5), Inches(7.0), Inches(12), Inches(0.8), safety_points, font_size=13)
    
    # ========== 第 2 页：应用场景与大模型工作流 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide2, prs)
    
    add_title_slide(slide2, "应用场景与大模型工作流")
    
    # 表格数据
    table_data = [
        ["场景", "大模型选择", "工作流", "价值"],
        ["智能设计方案生成", "Qwen2.5-72B（云端）", "输入需求→自动调用地勘数据→生成初步方案→工程师审核优化", "缩短设计周期 50%+"],
        ["水文数据智能分析", "GLM-Edge（边侧）", "实时监测数据→异常检测→趋势预测→预警推送", "提前 2-4 小时预警洪水风险"],
        ["工程文档自动审查", "Qwen2.5-Coder（云端）", "上传设计文档→规范合规性检查→问题标注→修改建议", "减少人工审查工作量 70%"],
        ["知识库问答助手", "本地微调小模型（端侧）", "自然语言提问→检索设计院知识库→精准回答", "新人培训效率提升 3 倍"]
    ]
    
    add_table(slide2, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), 5, 4, table_data)
    
    # ========== 第 3 页：端侧 Agent 部署 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide3, prs)
    
    add_title_slide(slide3, "端侧 Agent 部署 —— 智慧闸门闸控系统（辅助决策模式）")
    add_flow_diagram(slide3)
    
    # 表格数据
    table3_data = [
        ["场景", "端侧 Agent 功能", "人工介入", "价值"],
        ["洪水应急调度", "根据水位 + 降雨预测生成开闸建议方案", "✅ 必须确认", "缩短决策时间，责任清晰"],
        ["灌溉精准配水", "根据作物需水量 + 土壤湿度推荐配水计划", "✅ 必须确认", "减少经验依赖，优化资源配置"],
        ["设备故障诊断", "实时监测电机、传感器状态，提前预警", "⚠️ 可选确认", "预防性维护，减少停机"],
        ["异常告警", "发现数据异常自动推送告警 + 初步分析", "✅ 必须确认", "快速响应，避免误判"]
    ]
    
    add_table(slide3, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0), 5, 4, table3_data)
    
    # 核心原则
    principle_points = [
        "🤖 AI 负责：数据分析、方案生成、风险预警",
        "👤 人负责：最终决策、责任承担、例外处理",
        "🔒 系统保护：设置安全阈值，超出范围自动锁定需上级授权"
    ]
    add_bullet_points(slide3, Inches(0.5), Inches(4.0), Inches(12), Inches(1.0), principle_points, font_size=13)
    
    # ========== 第 4 页：客户侧延伸 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide4, prs)
    
    add_title_slide(slide4, "客户侧延伸 —— 水库/灌区的边端协同")
    
    # 左侧：边侧
    edge_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.5))
    edge_box.fill.solid()
    edge_box.fill.fore_color.rgb = LIGHT_GREEN
    edge_box.line.color.rgb = PRIMARY_COLOR
    edge_box.line.width = Pt(2)
    
    edge_text = "边侧（客户数据中心）\n\n"
    edge_text += "• 区域数据汇聚（多水库、多灌区）\n"
    edge_text += "• 本地大模型推理（无需上传云端）\n"
    edge_text += "• 与省院云端协同（模型更新、知识同步）"
    
    edge_frame = edge_box.text_frame
    edge_frame.clear()
    para = edge_frame.paragraphs[0]
    para.text = edge_text
    para.font.name = '微软雅黑'
    para.font.size = Pt(14)
    para.font.color.rgb = TEXT_COLOR
    para.alignment = PP_ALIGN.LEFT
    
    # 右侧：端侧
    device_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.2), Inches(5.8), Inches(2.5))
    device_box.fill.solid()
    device_box.fill.fore_color.rgb = LIGHT_ORANGE
    device_box.line.color.rgb = PRIMARY_COLOR
    device_box.line.width = Pt(2)
    
    device_text = "端侧（终端监测/控制设备）\n\n"
    device_text += "• 水位站、雨量站、流量计等监测设备嵌入 Agent\n"
    device_text += "• 闸门、泵站等控制设备智能决策\n"
    device_text += "• 巡检无人机/机器人自主巡护"
    
    device_frame = device_box.text_frame
    device_frame.clear()
    para = device_frame.paragraphs[0]
    para.text = device_text
    para.font.name = '微软雅黑'
    para.font.size = Pt(14)
    para.font.color.rgb = TEXT_COLOR
    para.alignment = PP_ALIGN.LEFT
    
    # 客户价值表格
    value_table_data = [
        ["价值维度", "具体收益"],
        ["管理效率", "减少人工巡检 60%，自动化报表生成"],
        ["决策质量", "数据驱动决策，减少经验依赖"],
        ["响应速度", "应急事件响应从小时级降至分钟级"],
        ["成本控制", "长期运维成本降低 30-40%"]
    ]
    
    add_table(slide4, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.0), 5, 2, value_table_data)
    
    # ========== 第 5 页：总结与展望 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide5, prs)
    
    add_title_slide(slide5, "总结与展望")
    
    # 第一阶段
    phase1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(3.8), Inches(1.8))
    phase1.fill.solid()
    phase1.fill.fore_color.rgb = LIGHT_BLUE
    phase1.line.color.rgb = PRIMARY_COLOR
    phase1.line.width = Pt(2)
    
    phase1_text = "第一阶段（2026 Q2-Q3）：搭建框架\n\n"
    phase1_text += "• 部署云端私有化大模型\n"
    phase1_text += "• 完成信息中心国产算力适配\n"
    phase1_text += "• 选择 1-2 个试点场景"
    
    phase1_frame = phase1.text_frame
    phase1_frame.clear()
    para = phase1_frame.paragraphs[0]
    para.text = phase1_text
    para.font.name = '微软雅黑'
    para.font.size = Pt(13)
    para.font.color.rgb = TEXT_COLOR
    para.alignment = PP_ALIGN.LEFT
    
    # 第二阶段
    phase2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.75), Inches(1.2), Inches(3.8), Inches(1.8))
    phase2.fill.solid()
    phase2.fill.fore_color.rgb = LIGHT_GREEN
    phase2.line.color.rgb = PRIMARY_COLOR
    phase2.line.width = Pt(2)
    
    phase2_text = "第二阶段（2026 Q4-2027 Q1）：场景落地\n\n"
    phase2_text += "• 智能设计方案生成上线\n"
    phase2_text += "• 智慧闸门端侧 Agent 试点\n"
    phase2_text += "• 培训设计人员使用新工具"
    
    phase2_frame = phase2.text_frame
    phase2_frame.clear()
    para = phase2_frame.paragraphs[0]
    para.text = phase2_text
    para.font.name = '微软雅黑'
    para.font.size = Pt(13)
    para.font.color.rgb = TEXT_COLOR
    para.alignment = PP_ALIGN.LEFT
    
    # 第三阶段
    phase3 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.2), Inches(3.8), Inches(1.8))
    phase3.fill.solid()
    phase3.fill.fore_color.rgb = LIGHT_ORANGE
    phase3.line.color.rgb = PRIMARY_COLOR
    phase3.line.width = Pt(2)
    
    phase3_text = "第三阶段（2027 Q2 起）：规模推广\n\n"
    phase3_text += "• 覆盖 80% 核心业务场景\n"
    phase3_text += "• 向客户（水库/灌区）推广边端方案\n"
    phase3_text += "• 形成行业标准与最佳实践"
    
    phase3_frame = phase3.text_frame
    phase3_frame.clear()
    para = phase3_frame.paragraphs[0]
    para.text = phase3_text
    para.font.name = '微软雅黑'
    para.font.size = Pt(13)
    para.font.color.rgb = TEXT_COLOR
    para.alignment = PP_ALIGN.LEFT
    
    # 愿景引用框
    vision_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.3), Inches(12.3), Inches(2.0))
    vision_box.fill.solid()
    vision_box.fill.fore_color.rgb = LIGHT_PURPLE
    vision_box.line.color.rgb = ACCENT_COLOR
    vision_box.line.width = Pt(3)
    
    vision_text = "让 AI Agent 成为水利设计师的智能助手，而非替代者\n安全、可控、高效地推动设计院数字化转型"
    
    vision_frame = vision_box.text_frame
    vision_frame.clear()
    para = vision_frame.paragraphs[0]
    para.text = vision_text
    para.font.name = '微软雅黑'
    para.font.size = Pt(18)
    para.font.color.rgb = PRIMARY_COLOR
    para.font.bold = True
    para.alignment = PP_ALIGN.CENTER
    para.space_after = Pt(10)
    
    # 保存文件
    output_file = "端云协同 AI Agent 在水利设计院的创新应用.pptx"
    prs.save(output_file)
    print("PPT generated successfully: " + output_file)
    return output_file

if __name__ == "__main__":
    create_ppt()
